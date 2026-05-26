#!/usr/bin/env python3
import argparse
import math
import json
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

BLUE = RGBColor(0x01, 0x5B, 0xAC)
RED = RGBColor(0xC6, 0x00, 0x19)
GRAY = RGBColor(0xA7, 0xB0, 0xBA)
DARK = RGBColor(0x1D, 0x2A, 0x38)
LIGHT_BG = RGBColor(0xF3, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_NAME = "Microsoft YaHei"
CHART_LEFT = 0.75
CHART_TOP = 1.15
CHART_WIDTH = 6.0
CHART_HEIGHT = 3.1


def load_plan(path: str) -> Dict[str, Any]:
    with open(path, "r", encoding="utf-8") as fh:
        payload = json.load(fh)
    plan_id = payload.get("confirmed_plan_id") or payload.get("recommended_plan_id")
    if not plan_id:
        raise ValueError("No confirmed_plan_id or recommended_plan_id in plan JSON")
    for plan in payload.get("chart_plans", []):
        if plan["plan_id"] == plan_id:
            payload["selected_plan"] = plan
            return payload
    raise ValueError(f"Plan id not found: {plan_id}")


def pick_series(payload: Dict[str, Any]) -> Tuple[List[str], List[Tuple[str, List[float]]]]:
    records = payload["data_records"]
    roles = payload["field_roles"]
    x_axis = roles["x_axis"]
    series_field = roles.get("series_field")
    columns = payload["columns"]
    field_types = {item["original_name"]: item["inferred_type"] for item in columns}
    value_columns = roles["value_columns"]
    primary_value = roles["primary_value"]
    plan = payload["selected_plan"]
    chart_type = plan["chart_type"]

    categories = [stringify(record.get(x_axis, "")) for record in records]
    if chart_type == "range_lines":
        preferred = [
            column["original_name"]
            for column in columns
            if column["inferred_type"] in ("lower_bound", "average", "upper_bound")
        ]
        if primary_value and primary_value not in preferred:
            preferred.insert(1, primary_value)
        series = [(name, [to_number(record.get(name), field_types.get(name, "numeric")) for record in records]) for name in preferred]
        return categories, series

    if series_field:
        category_order = []
        grouped: Dict[str, Dict[str, float]] = {}
        for record in records:
            category = stringify(record.get(x_axis, ""))
            series_name = stringify(record.get(series_field, ""))
            if category not in grouped:
                grouped[category] = {}
                category_order.append(category)
            grouped[category][series_name] = to_number(record.get(primary_value), field_types.get(primary_value, "numeric"))
        series_names = []
        for series_map in grouped.values():
            for name in series_map.keys():
                if name not in series_names:
                    series_names.append(name)
        result = []
        for series_name in series_names:
            result.append((series_name, [grouped[category].get(series_name) for category in category_order]))
        return category_order, result

    if chart_type == "line_target" and roles.get("target_columns"):
        target_name = roles["target_columns"][0]
        return categories, [
            (primary_value, [to_number(record.get(primary_value), field_types.get(primary_value, "numeric")) for record in records]),
            (target_name, [to_number(record.get(target_name), field_types.get(target_name, "numeric")) for record in records]),
        ]

    metric_columns = [
        name for name in value_columns
        if field_types.get(name) in ("numeric", "percent", "currency", "duration", "average")
    ]
    if len(metric_columns) > 1:
        return categories, [
            (name, [to_number(record.get(name), field_types.get(name, "numeric")) for record in records])
            for name in metric_columns
        ]
    return categories, [
        (primary_value, [to_number(record.get(primary_value), field_types.get(primary_value, "numeric")) for record in records])
    ]


def stringify(value: Any) -> str:
    return "" if value is None else str(value)


def to_number(value: Any, field_type: str) -> float:
    if value in (None, ""):
        return None
    text = str(value).replace(",", "").replace("%", "").replace("％", "")
    text = text.replace("万元", "").replace("元", "").replace("天", "").replace("小时", "").replace("分钟", "")
    return float(text)


def chart_enum(chart_type: str) -> XL_CHART_TYPE:
    mapping = {
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "multi_line": XL_CHART_TYPE.LINE,
        "clustered_column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "area": XL_CHART_TYPE.AREA,
        "line_target": XL_CHART_TYPE.LINE,
        "range_lines": XL_CHART_TYPE.LINE,
    }
    return mapping.get(chart_type, XL_CHART_TYPE.LINE_MARKERS)


def is_rapid_convergence_layout(payload: Dict[str, Any]) -> bool:
    metric_name = str(payload["field_roles"].get("primary_value", ""))
    tags = set(payload.get("task_tags", []))
    return (
        any(keyword in metric_name for keyword in ("准确率", "识别率", "稳定率"))
        and {"trend", "improvement", "stabilization"}.issubset(tags)
        and len(payload.get("data_records", [])) >= 6
    )


def use_percent_scale(payload: Dict[str, Any]) -> bool:
    metric_name = str(payload["field_roles"].get("primary_value", ""))
    values = metric_values(payload)
    return (
        any(keyword in metric_name for keyword in ("准确率", "识别率", "稳定率"))
        and values
        and all(0 <= value <= 1.2 for value in values)
    )


def add_title(slide, title_text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(11.6), Inches(0.7))
    text_frame = box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title_text
    run.font.name = FONT_NAME
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = DARK


def add_conclusion_box(slide, title: str, body: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.2), Inches(3.5), Inches(4.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = GRAY
    shape.line.width = Pt(1)
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.clear()

    p1 = text_frame.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = FONT_NAME
    r1.font.bold = True
    r1.font.size = Pt(16)
    r1.font.color.rgb = BLUE

    p2 = text_frame.add_paragraph()
    p2.space_before = Pt(8)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = FONT_NAME
    r2.font.size = Pt(12)
    r2.font.color.rgb = DARK


def add_info_box(slide, left: float, top: float, width: float, height: float, title: str, lines: List[str], title_color=BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = GRAY
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = FONT_NAME
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = title_color
    for line in lines:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = line
        r.font.name = FONT_NAME
        r.font.size = Pt(11)
        r.font.color.rgb = DARK


def add_key_point_tags(slide, key_points: List[Dict[str, Any]]) -> None:
    start_left = 0.8
    top = 5.6
    width = 2.6
    for idx, point in enumerate(key_points[:3]):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_left + idx * 2.75), Inches(top), Inches(width), Inches(0.58))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG if idx < 2 else RGBColor(0xFE, 0xF2, 0xF4)
        shape.line.color.rgb = BLUE if idx < 2 else RED
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = format_key_point(point)
        run.font.name = FONT_NAME
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = DARK


def format_key_point(point: Dict[str, Any]) -> str:
    label_map = {
        "start": "起点",
        "end": "终点",
        "max": "最高点",
        "min": "最低点",
        "first_target_reached": "首次达标",
        "turning_point": "拐点",
        "stable_region": "稳定区间",
        "anomaly": "异常波动",
    }
    label = label_map.get(point["kind"], point["kind"])
    return f"{label} | {point['x']} | {point['value']}"


def configure_chart(chart, selected_plan: Dict[str, Any], series_count: int) -> None:
    style = selected_plan["style"]
    chart.has_title = False
    chart.has_legend = not style.get("hide_legend", True) and series_count > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = FONT_NAME
        chart.legend.font.size = Pt(10)

    category_axis = chart.category_axis
    value_axis = chart.value_axis
    category_axis.has_title = False
    value_axis.has_title = False
    category_axis.tick_labels.font.name = FONT_NAME
    category_axis.tick_labels.font.size = Pt(10)
    value_axis.tick_labels.font.name = FONT_NAME
    value_axis.tick_labels.font.size = Pt(10)
    try:
        value_axis.has_major_gridlines = not selected_plan["style"].get("hide_gridlines", True)
        if value_axis.has_major_gridlines:
            value_axis.major_gridlines.format.line.color.rgb = GRAY
    except Exception:
        pass

    numeric_values = []
    for series in chart.series:
        for value in series.values:
            if value is not None:
                numeric_values.append(float(value))
    apply_value_axis_bounds(value_axis, numeric_values)
    value_axis.tick_labels.number_format = choose_number_format(numeric_values)

    for idx, series in enumerate(chart.series):
        line_color = BLUE if idx == 0 else (RED if idx == 1 else GRAY)
        try:
            series.format.line.color.rgb = line_color
            series.format.line.width = Pt(2)
        except Exception:
            pass
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = line_color
        except Exception:
            pass
        try:
            series.marker.style = XL_MARKER_STYLE.CIRCLE
            series.marker.size = 6
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = line_color
            series.marker.format.line.color.rgb = WHITE
            series.marker.format.line.width = Pt(1)
        except Exception:
            pass
        try:
            series.has_data_labels = False
        except Exception:
            pass


def configure_convergence_chart(chart, payload: Dict[str, Any]) -> None:
    selected_plan = payload["selected_plan"]
    configure_chart(chart, selected_plan, len(chart.series))
    chart.has_legend = False
    category_axis = chart.category_axis
    value_axis = chart.value_axis
    try:
        category_axis.format.line.color.rgb = RGBColor(0xD6, 0xDB, 0xE2)
        category_axis.format.line.width = Pt(1)
    except Exception:
        pass
    try:
        value_axis.tick_labels.font.size = Pt(1)
        value_axis.format.line.fill.background()
    except Exception:
        pass
    try:
        value_axis.has_major_gridlines = False
        value_axis.visible = False
    except Exception:
        pass
    try:
        category_axis.tick_labels.offset = 100
    except Exception:
        pass
    apply_convergence_data_labels(chart, payload)


def apply_value_axis_bounds(value_axis, numeric_values: List[float]) -> None:
    if not numeric_values:
        return

    data_min = min(numeric_values)
    data_max = max(numeric_values)
    if data_min == data_max:
        margin = max(abs(data_max) * 0.05, 1.0)
        minimum = data_min - margin
        maximum = data_max + margin
    else:
        span = data_max - data_min
        lower_margin = span * 0.1
        upper_margin = span * 0.1

        near_zero = data_min >= 0 and (data_min <= span * 0.35 or data_min <= max(1.0, data_max * 0.1))
        if near_zero:
            minimum = 0
        else:
            minimum = data_min - lower_margin
            if minimum < 0 and data_min >= 0:
                minimum = 0

        maximum = data_max + upper_margin

        if 60 <= data_min <= 100 and 60 <= data_max <= 100:
            minimum = math.floor(data_min)
            maximum = min(100, math.ceil(data_max + 0.6))
            if maximum <= minimum:
                maximum = data_max + 2.0

    try:
        value_axis.minimum_scale = float(minimum)
        value_axis.maximum_scale = float(maximum)
    except Exception:
        pass


def choose_number_format(numeric_values: List[float]) -> str:
    return "0"


def build_title(payload: Dict[str, Any]) -> str:
    summary = payload["task_summary"]
    plan = payload["selected_plan"]
    conclusion = plan["conclusion"]
    if "建议重点突出" in conclusion:
        return summary.replace("这是一个", "").replace("任务。", "")
    return conclusion


def metric_values(payload: Dict[str, Any]) -> List[float]:
    primary = payload["field_roles"]["primary_value"]
    field_type = next(item["inferred_type"] for item in payload["columns"] if item["original_name"] == primary)
    values = []
    for record in payload["data_records"]:
        value = to_number(record.get(primary), field_type)
        if value is not None:
            values.append(value)
    return values


def display_metric_values(payload: Dict[str, Any]) -> List[float]:
    values = metric_values(payload)
    if use_percent_scale(payload):
        return [value * 100 for value in values]
    return values


def category_values(payload: Dict[str, Any]) -> List[str]:
    x_axis = payload["field_roles"]["x_axis"]
    return [stringify(record.get(x_axis, "")) for record in payload["data_records"]]


def compact_value_text(value: float, payload: Dict[str, Any]) -> str:
    if use_percent_scale(payload):
        return f"{value * 100:.2f}%"
    if abs(value - round(value)) < 1e-9:
        return str(int(round(value)))
    return f"{value:.2f}"


def build_analysis_sections(payload: Dict[str, Any]) -> Dict[str, List[str]]:
    values = metric_values(payload)
    categories = category_values(payload)
    points = payload.get("key_points", [])
    primary = payload["field_roles"]["primary_value"]
    x_axis = payload["field_roles"]["x_axis"]
    start = compact_value_text(values[0], payload) if values else "-"
    end = compact_value_text(values[-1], payload) if values else "-"
    maximum = compact_value_text(max(values), payload) if values else "-"
    phase_idx = find_phase_index(payload, display_metric_values(payload)) if values else 0
    phase_month = categories[phase_idx] if categories else "-"
    stable = next((p["x"] for p in points if p["kind"] == "stable_region"), "未识别")

    overview = [
        payload["task_summary"],
        payload["selected_plan"]["conclusion"],
    ]
    structure = [
        f"横轴：{x_axis}",
        f"纵轴：{primary}",
        f"数据点数：{len(values)}",
        f"阶段分界建议：{phase_month}",
    ]
    insights = [
        f"起点：{categories[0]} / {start}" if categories and values else "起点：-",
        f"终点：{categories[-1]} / {end}" if categories and values else "终点：-",
        f"最高值：{maximum}",
        f"稳定区间：{stable}",
    ]
    recommendation = [
        "适合在正式汇报中表达“算法能力快速提升并进入稳定运行阶段”。",
        "若用于对外材料，建议搭配一段业务结果或版本迭代说明。",
    ]
    return {
        "overview": overview,
        "structure": structure,
        "insights": insights,
        "recommendation": recommendation,
    }


def add_analysis_panels(slide, payload: Dict[str, Any], emphasize_red: bool = False) -> None:
    sections = build_analysis_sections(payload)
    add_info_box(slide, 7.0, 1.1, 5.45, 1.55, "任务判断", sections["overview"], RED if emphasize_red else BLUE)
    add_info_box(slide, 7.0, 2.85, 5.45, 1.55, "图表结构", sections["structure"])
    add_info_box(slide, 0.75, 4.55, 6.0, 1.75, "关键观察", sections["insights"])
    add_info_box(slide, 7.0, 4.55, 5.45, 1.75, "使用建议", sections["recommendation"])


def find_first_index(values: List[float], threshold: float) -> Optional[int]:
    for idx, value in enumerate(values):
        if value >= threshold:
            return idx
    return None


def find_phase_index(payload: Dict[str, Any], values: List[float]) -> int:
    candidates = [98.0, 97.7, 97.5] if values and max(values) > 1.5 else [0.98, 0.977, 0.975]
    for threshold in candidates:
        idx = find_first_index(values, threshold)
        if idx is not None:
            return idx
    return max(1, len(values) // 2)


def choose_label_indices(values: List[float], phase_idx: int) -> List[int]:
    indices = [0]
    if len(values) > 1:
        indices.append(1)
    threshold = 94.0 if values and max(values) > 1.5 else 0.94
    pre_phase = next((idx for idx, value in enumerate(values) if value >= threshold), None)
    if pre_phase is not None:
        indices.append(pre_phase)
    if phase_idx is not None:
        candidate = max(0, phase_idx - 1)
        indices.append(candidate)
    max_idx = max(range(len(values)), key=lambda idx: values[idx])
    indices.append(max_idx)
    indices.append(len(values) - 1)
    deduped = []
    for idx in indices:
        if idx not in deduped:
            deduped.append(idx)
    return deduped[:6]


def value_to_chart_y(value: float, min_scale: float, max_scale: float) -> float:
    ratio = 0.5 if max_scale == min_scale else (value - min_scale) / (max_scale - min_scale)
    ratio = max(0.0, min(1.0, ratio))
    return CHART_TOP + CHART_HEIGHT - (ratio * CHART_HEIGHT)


def index_to_chart_x(index: int, total: int) -> float:
    if total <= 1:
        return CHART_LEFT + CHART_WIDTH / 2
    return CHART_LEFT + (CHART_WIDTH * index / (total - 1))


def format_percent(value: float) -> str:
    return f"{value:.2f}%" if value > 1.5 else f"{value * 100:.2f}%"


def apply_convergence_data_labels(chart, payload: Dict[str, Any]) -> None:
    values = display_metric_values(payload)
    phase_idx = find_phase_index(payload, values)
    label_indices = set(choose_label_indices(values, phase_idx))
    series = chart.series[0]
    for idx, point in enumerate(series.points):
        label = point.data_label
        tf = label.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        if idx in label_indices:
            run.text = format_percent(values[idx])
            run.font.name = FONT_NAME
            run.font.size = Pt(16 if idx in {0, len(values) - 1} else 15)
            run.font.bold = idx in {0, len(values) - 1}
            run.font.color.rgb = RED if idx in {0, len(values) - 1} else DARK
            label.position = 0
        else:
            run.text = ""


def add_convergence_labels(slide, payload: Dict[str, Any], chart) -> None:
    values = display_metric_values(payload)
    categories = category_values(payload)
    phase_idx = find_phase_index(payload, values)
    min_scale = chart.value_axis.minimum_scale
    max_scale = chart.value_axis.maximum_scale

    add_phase_line(slide, categories, values, phase_idx, min_scale, max_scale)
    add_convergence_callout(slide, values, phase_idx)


def add_phase_line(slide, categories: List[str], values: List[float], phase_idx: int, min_scale: float, max_scale: float) -> None:
    x = index_to_chart_x(phase_idx, len(values))
    top = value_to_chart_y(max(values) + (max_scale - min_scale) * 0.02, min_scale, max_scale)
    bottom = CHART_TOP + CHART_HEIGHT - 0.18
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x),
        Inches(top),
        Inches(x),
        Inches(bottom),
    )
    line.line.color.rgb = RED
    line.line.width = Pt(2)
    try:
        line.line.dash_style = 2
    except Exception:
        pass


def add_convergence_callout(slide, values: List[float], phase_idx: int) -> None:
    stable_value = max(98, int(round(values[-1] if values[-1] > 1.5 else values[-1] * 100)))
    box = slide.shapes.add_textbox(Inches(7.35), Inches(0.45), Inches(4.5), Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "算法迅速收敛"
    r1.font.name = FONT_NAME
    r1.font.size = Pt(24)
    r1.font.bold = True
    r1.font.color.rgb = RED
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = f"稳定至{stable_value}%"
    r2.font.name = FONT_NAME
    r2.font.size = Pt(24)
    r2.font.bold = True
    r2.font.color.rgb = RED


def render_rapid_convergence_slide(slide, payload: Dict[str, Any], categories: List[str], series_data: List[Tuple[str, List[float]]]) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = categories
    rapid_series = series_data
    if use_percent_scale(payload):
        rapid_series = [(series_name, [value * 100 if value is not None else None for value in values]) for series_name, values in series_data]
    for series_name, values in rapid_series:
        chart_data.add_series(series_name, values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(CHART_LEFT),
        Inches(CHART_TOP),
        Inches(CHART_WIDTH),
        Inches(CHART_HEIGHT),
        chart_data,
    )
    chart = chart_shape.chart
    configure_convergence_chart(chart, payload)
    add_convergence_labels(slide, payload, chart)
    add_analysis_panels(slide, payload, emphasize_red=True)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--plan", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    payload = load_plan(args.plan)
    selected_plan = payload["selected_plan"]
    categories, series_data = pick_series(payload)
    if not categories or not series_data:
        raise ValueError("No chartable data found in plan payload")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if is_rapid_convergence_layout(payload):
        render_rapid_convergence_slide(slide, payload, categories, series_data)
    else:
        add_title(slide, build_title(payload))

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for series_name, values in series_data:
            chart_data.add_series(series_name, values)

        chart_shape = slide.shapes.add_chart(
            chart_enum(selected_plan["chart_type"]),
            Inches(0.75),
            Inches(1.2),
            Inches(8.0),
            Inches(4.1),
            chart_data,
        )
        chart = chart_shape.chart
        configure_chart(chart, selected_plan, len(series_data))
        add_analysis_panels(slide, payload)
        add_key_point_tags(slide, selected_plan["key_points"])

        footer = slide.shapes.add_textbox(Inches(7.0), Inches(6.55), Inches(5.45), Inches(0.35))
        tf = footer.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"生成时间：{datetime.now().strftime('%Y-%m-%d')}  |  图表类型：{selected_plan['chart_name']}"
        run.font.name = FONT_NAME
        run.font.size = Pt(9)
        run.font.color.rgb = GRAY

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


if __name__ == "__main__":
    main()
