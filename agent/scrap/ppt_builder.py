"""镔鑫废钢检判 · 单页 PPT 汇报页 自研构建器（不依赖同事 skill）。

设计目标（解决同事 skill 默认产物的痛点）：
  1. 单线主图（聚焦主料识别率），不再多线杂乱
  2. 显式目标线（默认 95%）+ 4 个关键节点标签（起/终/最高/最低）
  3. 元信息卡：数据来源 / 周期 / 有效车次 / 错判车次（让领导一眼知道"这数据怎么来的"）
  4. KPI 大数字 + 达标徽章 + 一句话结论（直接给观点）
  5. 错判 Top 5 真实车牌（给案例）
  6. 改进建议 3 条（动态根据数据生成）

布局（13.33 x 7.5 英寸 16:9）：
  ┌──────────────────────────────────────────────────────────────────┐
  │ 顶部品牌带（深蓝） 0~0.6 in                                       │
  ├──────────────────────────────────────────────────────────────────┤
  │ 大标题 0.7~1.4 in                                                 │
  ├──────────────────────────────────────────────────────────────────┤
  │ 元信息卡 1.5~2.2 in（4 等分横排 RoundRect）                       │
  ├──────────────────────────────────────────────────────────────────┤
  │ 主图区 2.4~5.3 / 0.4~7.8 │ KPI 卡区 2.4~5.3 / 8.0~12.93           │
  ├──────────────────────────────────────────────────────────────────┤
  │ 错判 Top5 5.5~7.0 / 0.4~7.8 │ 改进建议 5.5~7.0 / 8.0~12.93        │
  ├──────────────────────────────────────────────────────────────────┤
  │ 页脚 7.1~7.4                                                      │
  └──────────────────────────────────────────────────────────────────┘
"""
from __future__ import annotations

import logging
from datetime import datetime
from pathlib import Path
from typing import Iterable, List, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

from agent.scrap.models import DailyScrapStats, TruckStat

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# 视觉规范（中冶赛迪 + 配色统一）
# ---------------------------------------------------------------------------
BRAND_BLUE = RGBColor(0x01, 0x5B, 0xAC)  # 主蓝
BRAND_DEEP = RGBColor(0x0A, 0x36, 0x6E)  # 深蓝（顶部带）
ACCENT_RED = RGBColor(0xC6, 0x00, 0x19)  # 关键点 + 目标线
SUCCESS_GREEN = RGBColor(0x2E, 0x7D, 0x32)  # 达标
WARN_ORANGE = RGBColor(0xE6, 0x7E, 0x22)  # 接近
DANGER_RED = RGBColor(0xC0, 0x39, 0x2B)  # 远低于
GRAY_DARK = RGBColor(0x2D, 0x3A, 0x4A)
GRAY_MID = RGBColor(0x6B, 0x77, 0x83)
GRAY_LIGHT = RGBColor(0xC9, 0xD3, 0xDC)
BG_LIGHT = RGBColor(0xF3, 0xF7, 0xFB)
BG_INFO = RGBColor(0xE8, 0xF1, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_FAMILY = "Microsoft YaHei"


# ---------------------------------------------------------------------------
# 工具：常用样式
# ---------------------------------------------------------------------------
def _set_text(
    paragraph,
    text: str,
    *,
    size: int = 11,
    bold: bool = False,
    color: RGBColor = GRAY_DARK,
    align=None,
) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align


def _add_textbox(
    slide, left: float, top: float, width: float, height: float
):
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )


def _add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor,
    line: Optional[RGBColor] = None,
    line_width_pt: float = 0.5,
    rounded: bool = True,
):
    shape_type = (
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    )
    shp = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width_pt)
    shp.shadow.inherit = False
    return shp


# ---------------------------------------------------------------------------
# 顶部品牌带 + 标题
# ---------------------------------------------------------------------------
def _draw_brand_strip(slide, cycle_label: str) -> None:
    """页面最顶端的深蓝色品牌带（13.33 宽 × 0.55 高）"""
    _add_rect(
        slide, 0, 0, 13.333, 0.55, fill=BRAND_DEEP, rounded=False
    )
    # 左：BriefMe｜中冶赛迪
    box = _add_textbox(slide, 0.4, 0.06, 6.0, 0.45)
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_text(p, "BriefMe", size=18, bold=True, color=WHITE)
    _set_text(p, "  ｜  中冶赛迪（重庆）信息技术有限公司", size=11, color=GRAY_LIGHT)
    # 右：周期
    box = _add_textbox(slide, 7.0, 0.10, 6.0, 0.40)
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _set_text(p, f"统计周期：{cycle_label}", size=12, bold=True, color=WHITE)


def _draw_title(slide) -> None:
    """大标题 + 副标题"""
    box = _add_textbox(slide, 0.4, 0.7, 12.93, 0.75)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    _set_text(
        p1, "镔鑫废钢检判 · 主料识别率周期分析",
        size=22, bold=True, color=BRAND_DEEP,
    )
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(2)
    _set_text(
        p2,
        "基于赛迪 AI 视觉检判结果 vs 现场人工质检 1 人对照",
        size=11, color=GRAY_MID,
    )


# ---------------------------------------------------------------------------
# 元信息卡（4 个并排 RoundRect）
# ---------------------------------------------------------------------------
def _draw_metainfo_cards(
    slide,
    *,
    source_label: str,
    cycle_label: str,
    eligible_total: int,
    error_total: int,
) -> None:
    """4 个并排卡片：数据来源 / 周期 / 有效车次 / 错判车次"""
    # 数据来源单独缩短（避免文字溢出卡片）
    short_source = (
        source_label.split("（")[0]
        if "（" in source_label
        else source_label
    )
    cards = [
        ("数据来源", short_source, BRAND_BLUE),
        ("统计周期", cycle_label, BRAND_BLUE),
        ("有效车次", f"{eligible_total} 辆", SUCCESS_GREEN),
        ("错判车次", f"{error_total} 辆", DANGER_RED),
    ]
    top = 1.55
    height = 0.65
    total_w = 12.93 - 0.4  # 12.53
    gap = 0.12
    card_w = (total_w - gap * 3) / 4
    left = 0.4
    for label, value, accent in cards:
        _add_rect(
            slide, left, top, card_w, height,
            fill=BG_INFO, line=GRAY_LIGHT, line_width_pt=0.5,
        )
        # 左侧 4px 强调条
        _add_rect(
            slide, left, top, 0.06, height,
            fill=accent, rounded=False,
        )
        # 文本
        tb = _add_textbox(slide, left + 0.18, top + 0.05, card_w - 0.22, height - 0.1)
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        p1 = tf.paragraphs[0]
        _set_text(p1, label, size=10, color=GRAY_MID)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        _set_text(p2, value, size=14, bold=True, color=accent)
        left += card_w + gap


# ---------------------------------------------------------------------------
# 主图：单线主料识别率 + 目标线 + 关键节点标签
# ---------------------------------------------------------------------------
def _disable_marker(series) -> None:
    """python-pptx 没有直接 marker.style=NONE，用 lxml 注入 <c:marker><c:symbol val='none'/></c:marker>"""
    ser = series._element
    nsmap = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
    # 删除既有 marker
    for old in ser.findall(qn("c:marker")):
        ser.remove(old)
    marker = etree.SubElement(ser, qn("c:marker"))
    symbol = etree.SubElement(marker, qn("c:symbol"))
    symbol.set("val", "none")


def _set_dash_style(series, dash: str = "dash") -> None:
    """给折线设虚线风格。dash 可选: dash / dashDot / lgDash / sysDash 等"""
    ser = series._element
    spPr = ser.find(qn("c:spPr"))
    if spPr is None:
        spPr = etree.SubElement(ser, qn("c:spPr"))
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(
            spPr,
            qn("a:ln"),
        )
    # 删除原 prstDash
    for old in ln.findall(qn("a:prstDash")):
        ln.remove(old)
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", dash)


def _draw_chart(
    slide,
    dates: List[str],
    rates: List[float],
    target_pct: float,
) -> None:
    """画主料识别率折线 + 红色虚线目标线，返回 chart 对象"""
    left, top, width, height = 0.4, 2.35, 7.6, 3.0

    # 标题文本（放在图上方，不占图内空间）
    tb = _add_textbox(slide, left, top - 0.05, width, 0.35)
    p = tb.text_frame.paragraphs[0]
    _set_text(p, "主料识别率周期趋势", size=13, bold=True, color=BRAND_DEEP)
    p2 = tb.text_frame.add_paragraph()
    _set_text(
        p2, "（红色虚线为目标值）", size=9, color=GRAY_MID,
    )

    # 横轴日期格式化为 MM-DD
    short_dates = [_short_date(d) for d in dates]

    chart_data = CategoryChartData()
    chart_data.categories = short_dates
    chart_data.add_series("主料识别率(%)", rates)
    chart_data.add_series(
        f"目标值 {target_pct:.0f}%", [target_pct] * len(rates)
    )

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(left),
        Inches(top + 0.35),
        Inches(width),
        Inches(height - 0.4),
        chart_data,
    )
    chart = chart_shape.chart

    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = FONT_FAMILY
    chart.legend.font.size = Pt(9)

    # 系列 1 = 主料识别率，蓝色实线 + 圆形 marker
    main_series = chart.series[0]
    main_series.format.line.color.rgb = BRAND_BLUE
    main_series.format.line.width = Pt(2.25)
    # marker 颜色
    try:
        main_series.marker.format.fill.solid()
        main_series.marker.format.fill.fore_color.rgb = BRAND_BLUE
        main_series.marker.format.line.color.rgb = WHITE
        main_series.marker.format.line.width = Pt(1)
    except Exception:
        pass

    # 系列 2 = 目标线，红色虚线，无 marker
    target_series = chart.series[1]
    target_series.format.line.color.rgb = ACCENT_RED
    target_series.format.line.width = Pt(1.5)
    _set_dash_style(target_series, "dash")
    _disable_marker(target_series)

    # Y 轴
    val_axis = chart.value_axis
    val_axis.has_title = True
    val_axis.axis_title.text_frame.text = "主料识别率（%）"
    for p in val_axis.axis_title.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = FONT_FAMILY
            r.font.size = Pt(10)
            r.font.color.rgb = GRAY_DARK
    val_axis.tick_labels.font.name = FONT_FAMILY
    val_axis.tick_labels.font.size = Pt(9)
    val_axis.minimum_scale = 0
    val_axis.maximum_scale = 100
    val_axis.major_unit = 20
    val_axis.tick_labels.number_format = "0"

    # X 轴
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.name = FONT_FAMILY
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.has_title = False

    # 数据标签：序列级整体启用，整数百分比，小字号
    # 这样多个相同数值（如连续 100%）视觉叠成一个，避免重叠
    dlabels = main_series.data_labels
    dlabels.show_value = True
    dlabels.number_format = '0"%"'
    dlabels.position = XL_LABEL_POSITION.ABOVE
    dlabels.font.name = FONT_FAMILY
    dlabels.font.size = Pt(9)
    dlabels.font.bold = True
    dlabels.font.color.rgb = BRAND_DEEP

    # 目标线不操作 .data_labels（python-pptx 默认不会创建标签节点）

    # 高亮极值点：min/max 用红色 bold 大字号 override
    if rates:
        n = len(rates)
        max_idx = max(range(n), key=lambda i: rates[i])
        min_idx = min(range(n), key=lambda i: rates[i])
        for idx, pos in [
            (max_idx, XL_LABEL_POSITION.ABOVE),
            (min_idx, XL_LABEL_POSITION.BELOW),
        ]:
            try:
                lbl = main_series.points[idx].data_label
                lbl.position = pos
                lbl.font.size = Pt(11)
                lbl.font.bold = True
                lbl.font.color.rgb = ACCENT_RED
            except Exception:
                pass


# ---------------------------------------------------------------------------
# KPI 卡区
# ---------------------------------------------------------------------------
def _draw_kpi_panel(
    slide,
    *,
    overall_rate_pct: float,
    target_pct: float,
    eligible_total: int,
    error_total: int,
    days_meet_target: int,
    total_days: int,
) -> None:
    left, top, width, height = 8.2, 2.35, 4.93, 3.0

    # 外框
    _add_rect(
        slide, left, top, width, height,
        fill=BG_LIGHT, line=GRAY_LIGHT, line_width_pt=0.75,
    )

    # 标题
    tb = _add_textbox(slide, left + 0.25, top + 0.1, width - 0.5, 0.4)
    p = tb.text_frame.paragraphs[0]
    _set_text(p, "周期核心 KPI", size=13, bold=True, color=BRAND_DEEP)

    # 大数字
    big_box = _add_textbox(slide, left + 0.25, top + 0.55, width - 0.5, 1.2)
    tf = big_box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_text(
        p, f"{overall_rate_pct:.2f}%",
        size=44, bold=True,
        color=_status_color(overall_rate_pct, target_pct),
    )
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)
    _set_text(
        p2, f"主料识别率（目标 {target_pct:.0f}%）",
        size=11, color=GRAY_MID,
    )

    # 达标徽章
    badge_text, badge_color = _status_badge(overall_rate_pct, target_pct)
    badge_left = left + (width - 1.4) / 2
    _add_rect(
        slide, badge_left, top + 1.85, 1.4, 0.32,
        fill=badge_color, rounded=True,
    )
    badge_box = _add_textbox(slide, badge_left, top + 1.87, 1.4, 0.28)
    p = badge_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_text(p, badge_text, size=11, bold=True, color=WHITE)

    # 二级指标（2 列 × 2 行）
    sub_top = top + 2.30
    grid = [
        (f"{error_total}", "错判车次", DANGER_RED),
        (f"{eligible_total}", "有效车次", BRAND_BLUE),
        (f"{days_meet_target}/{total_days}", "达标天数", SUCCESS_GREEN),
        (
            f"{(error_total / eligible_total * 100.0):.1f}%"
            if eligible_total > 0 else "—",
            "错判占比",
            WARN_ORANGE,
        ),
    ]
    cell_w = (width - 0.4) / 2
    cell_h = 0.32
    for i, (val, lab, c) in enumerate(grid):
        r = i // 2
        col = i % 2
        x = left + 0.2 + col * cell_w
        y = sub_top + r * cell_h
        tb2 = _add_textbox(slide, x, y, cell_w, cell_h)
        p = tb2.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _set_text(p, val, size=12, bold=True, color=c)
        _set_text(p, f"  {lab}", size=10, color=GRAY_MID)


def _status_color(value_pct: float, target_pct: float) -> RGBColor:
    if value_pct >= target_pct:
        return SUCCESS_GREEN
    if value_pct >= target_pct - 10:
        return WARN_ORANGE
    return DANGER_RED


def _status_badge(value_pct: float, target_pct: float) -> tuple:
    if value_pct >= target_pct:
        return ("已达标", SUCCESS_GREEN)
    if value_pct >= target_pct - 10:
        return ("接近目标", WARN_ORANGE)
    return ("远低于目标", DANGER_RED)


# ---------------------------------------------------------------------------
# 错判 Top 5 表格
# ---------------------------------------------------------------------------
def _draw_topn_errors(
    slide, top_errors: List[TruckStat]
) -> None:
    # 上移 0.10 + 增高 0.20，给 6 行表格留足空间
    left, top, width, height = 0.4, 5.40, 7.6, 1.85

    # 标题
    tb = _add_textbox(slide, left, top - 0.05, width, 0.32)
    p = tb.text_frame.paragraphs[0]
    _set_text(
        p, f"错判车次 Top {len(top_errors)}（按差异率降序）",
        size=12, bold=True, color=BRAND_DEEP,
    )

    table_top = top + 0.28
    table_h = height - 0.28

    if not top_errors:
        # 空态
        _add_rect(
            slide, left, table_top, width, table_h,
            fill=BG_LIGHT, line=GRAY_LIGHT,
        )
        tb = _add_textbox(slide, left, table_top, width, table_h)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_text(
            p, "本周期无错判车次（恭喜，全部主料判定一致）",
            size=11, color=SUCCESS_GREEN, bold=True,
        )
        return

    rows = len(top_errors) + 1
    cols = 5
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(table_top),
        Inches(width), Inches(table_h),
    )
    tbl = table_shape.table

    # 列宽（in）
    widths = [1.0, 1.4, 0.8, 2.2, 2.2]
    total = sum(widths)
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(width * w / total)

    # 行高：表头略高，数据行紧凑
    header_h = 0.30
    row_h = (table_h - header_h) / max(len(top_errors), 1)
    tbl.rows[0].height = Inches(header_h)
    for r in range(1, rows):
        tbl.rows[r].height = Inches(row_h)

    headers = ["日期", "车牌", "工位", "AI 主料 → 占比", "人工 主料 → 占比"]
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_DEEP
        cell.text_frame.clear()
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_text(p, h, size=10, bold=True, color=WHITE)

    for r, t in enumerate(top_errors, start=1):
        row_fill = BG_LIGHT if r % 2 == 1 else WHITE
        ai = _fmt_main(t.ai_main, with_rate=True)
        manual = _fmt_main(t.manual_main, with_rate=True)
        cells = [t.date, t.car_number, str(t.station_number), ai, manual]
        for c, txt in enumerate(cells):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_fill
            cell.text_frame.clear()
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _set_text(p, txt, size=9, color=GRAY_DARK)


def _fmt_main(entry, *, with_rate: bool = False) -> str:
    if entry is None:
        return "—"
    from agent.scrap.dict import get_material_name

    name = get_material_name(entry.steel_type, entry.steel_level)
    if not with_rate or entry.rate is None:
        return name
    return f"{name}  {entry.rate:.1f}%"


# ---------------------------------------------------------------------------
# 改进建议
# ---------------------------------------------------------------------------
def _draw_recommendations(
    slide, tips: List[str]
) -> None:
    # 与错判表对齐：5.40 起，1.85 高
    left, top, width, height = 8.2, 5.40, 4.93, 1.85

    # 外框
    _add_rect(
        slide, left, top, width, height,
        fill=BG_INFO, line=BRAND_BLUE, line_width_pt=0.75,
    )

    # 标题
    tb = _add_textbox(slide, left + 0.2, top + 0.1, width - 0.4, 0.4)
    p = tb.text_frame.paragraphs[0]
    _set_text(p, "改进建议", size=12, bold=True, color=BRAND_DEEP)

    # bullets
    body = _add_textbox(slide, left + 0.2, top + 0.5, width - 0.4, height - 0.55)
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, tip in enumerate(tips):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        _set_text(p, "• ", size=10, bold=True, color=BRAND_BLUE)
        _set_text(p, tip, size=10, color=GRAY_DARK)


# ---------------------------------------------------------------------------
# 页脚
# ---------------------------------------------------------------------------
def _draw_footer(slide, source_label: str) -> None:
    left, top, width, height = 0.4, 7.18, 12.93, 0.27
    tb = _add_textbox(slide, left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    _set_text(
        p,
        f"自动生成 · BriefMe / 中冶赛迪    |    数据源：{source_label}    |    生成时间：{now}",
        size=8, color=GRAY_MID,
    )


# ---------------------------------------------------------------------------
# 工具函数
# ---------------------------------------------------------------------------
def _short_date(d: str) -> str:
    """YYYY-MM-DD → MM-DD"""
    parts = d.split("-")
    if len(parts) == 3:
        return f"{int(parts[1]):02d}-{int(parts[2]):02d}"
    return d


def _format_cycle(start: str, end: str) -> str:
    if start == end:
        return start
    return f"{start} ~ {end}"


def _generate_tips(
    daily_records: List[tuple],
    top_errors: List[TruckStat],
    target_pct: float,
    overall_rate: float,
) -> List[str]:
    """根据数据动态生成 3 条改进建议。

    daily_records: List[(date, eligible, errors, accuracy_pct)]
    """
    tips: List[str] = []
    if not daily_records:
        return ["本周期暂无有效数据，建议扩大时间范围或检查数据来源。"]

    # 1. 达标天数 / 最差日 / 错判最多日
    meet = [r for r in daily_records if r[3] is not None and r[3] >= target_pct]
    not_none = [r for r in daily_records if r[3] is not None]
    if not_none:
        worst = min(not_none, key=lambda r: r[3])
        most_err = max(not_none, key=lambda r: r[2])
        tips.append(
            f"周期内 {len(meet)}/{len(daily_records)} 天达标（≥{target_pct:.0f}%），"
            f"最低识别率出现在 {worst[0]}（{worst[3]:.1f}%）。"
        )
        if most_err[2] > 0:
            tips.append(
                f"错判最集中在 {most_err[0]}，当天有 {most_err[2]} 辆判错；"
                "建议优先排查该日现场情况（光照/装载/料堆形态）。"
            )

    # 2. 错判主料类型集中度
    if top_errors:
        from collections import Counter
        from agent.scrap.dict import get_material_name

        ai_types = Counter(
            get_material_name(t.ai_main.steel_type, t.ai_main.steel_level)
            for t in top_errors
            if t.ai_main is not None
        )
        if ai_types:
            most_common, count = ai_types.most_common(1)[0]
            tips.append(
                f"AI 错判最集中的料型是「{most_common}」({count} 次)；"
                "建议补充该料型在不同光照/角度下的样本，下次模型迭代覆盖。"
            )

    # 3. 整体差距
    gap = target_pct - overall_rate
    if gap > 20:
        tips.append(
            f"整体识别率距目标 {target_pct:.0f}% 仍有 {gap:.1f} 个百分点的差距，"
            "建议把当前阶段目标暂调至更接近的中期值（如 70~80%）。"
        )

    # 截到 3 条
    if not tips:
        tips = ["数据指标稳定，建议继续保持现有作业规范。"]
    return tips[:3]


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------
def build_binxin_ppt(
    stats_list: List[DailyScrapStats],
    save_path: Path,
    *,
    start_date: str,
    end_date: str,
    source_label: str = "镔鑫废钢检判系统（赛迪 AI vs 人工质检）",
    target_pct: float = 95.0,
    top_n_errors: int = 5,
) -> Path:
    """构建一页镔鑫专属 PPT 汇报页。

    Args:
        stats_list: 每日统计列表（按日期升序）
        save_path: 输出 .pptx 路径
        start_date / end_date: 周期边界（仅用于显示）
        source_label: 数据源标签（页眉/页脚/元信息卡）
        target_pct: 主料识别率目标值（默认 95%）
        top_n_errors: 错判 Top N 显示几行

    Returns:
        save_path
    """
    # 过滤掉 eligible=0 的天
    valid_days = [s for s in stats_list if s.eligible_trucks > 0]

    daily_records = []  # (date, eligible, errors, accuracy_pct)
    dates: List[str] = []
    rates: List[float] = []
    for s in valid_days:
        eligible = s.eligible_trucks
        errors = eligible - s.main_same_count
        acc = s.accuracy_rate or 0.0
        daily_records.append((s.date, eligible, errors, acc))
        dates.append(s.date)
        rates.append(acc)

    eligible_total = sum(d[1] for d in daily_records)
    error_total = sum(d[2] for d in daily_records)
    days_meet_target = sum(
        1 for d in daily_records if d[3] is not None and d[3] >= target_pct
    )
    overall_rate = (
        (eligible_total - error_total) / eligible_total * 100.0
        if eligible_total > 0 else 0.0
    )

    # 错判 Top N（跨所有日子，按 diff_rate 降序；diff_rate 为 None 的放最后）
    all_errors: List[TruckStat] = []
    for s in valid_days:
        for t in s.trucks:
            if t.main_same is False:
                all_errors.append(t)
    all_errors.sort(
        key=lambda t: (t.diff_rate is None, -(t.diff_rate or 0.0))
    )
    top_errors = all_errors[:top_n_errors]

    tips = _generate_tips(daily_records, top_errors, target_pct, overall_rate)
    cycle_label = _format_cycle(start_date, end_date)

    # ------ 构建 PPT ------
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    _draw_brand_strip(slide, cycle_label)
    _draw_title(slide)
    _draw_metainfo_cards(
        slide,
        source_label=source_label,
        cycle_label=cycle_label,
        eligible_total=eligible_total,
        error_total=error_total,
    )
    _draw_chart(slide, dates, rates, target_pct)
    _draw_kpi_panel(
        slide,
        overall_rate_pct=overall_rate,
        target_pct=target_pct,
        eligible_total=eligible_total,
        error_total=error_total,
        days_meet_target=days_meet_target,
        total_days=len(daily_records),
    )
    _draw_topn_errors(slide, top_errors)
    _draw_recommendations(slide, tips)
    _draw_footer(slide, source_label)

    save_path = Path(save_path).resolve()
    save_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(save_path))
    logger.info("镔鑫专属 PPT 已生成: %s", save_path)
    return save_path
