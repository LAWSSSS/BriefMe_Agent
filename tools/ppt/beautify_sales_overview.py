"""美化《销售数智化一期功能总图》PPT 单页架构图。

设计原则：
- 不动整体结构与卡片相对位置；
- 仅放大字号、同步微调文本框高度避免溢出、个别颜色细节优化；
- 通过 shape index 精确定位（已与源文件 shape 列表对齐，迁移前请重新核对）。

用法：
    python tools/ppt/beautify_sales_overview.py \
        --src 演示文件/销售数智化一期功能总图.pptx \
        --dst 演示文件/销售数智化一期功能总图.美化版.pptx
"""

from __future__ import annotations

import argparse
import logging
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Cm, Pt

logger = logging.getLogger("beautify_sales_overview")


# ---------------------------------------------------------------------------
# 通用工具
# ---------------------------------------------------------------------------


def _set_run_size(shape, size_pt: float) -> None:
    """把 text_frame 内所有 run 的字号统一设置为 size_pt。"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size_pt)


def _set_run_color(shape, hex_rgb: str) -> None:
    if not shape.has_text_frame:
        return
    rgb = RGBColor.from_string(hex_rgb)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = rgb


def _set_geometry(shape, *, top_cm: float | None = None, height_cm: float | None = None,
                  left_cm: float | None = None, width_cm: float | None = None) -> None:
    if top_cm is not None:
        shape.top = Cm(top_cm)
    if height_cm is not None:
        shape.height = Cm(height_cm)
    if left_cm is not None:
        shape.left = Cm(left_cm)
    if width_cm is not None:
        shape.width = Cm(width_cm)


def _vertical_center(shape) -> None:
    """让 text_frame 内文字垂直居中，避免大字偏上。"""
    if shape.has_text_frame:
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _set_margins(shape, *, left=0.1, right=0.1, top=0.05, bottom=0.05) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.margin_left = Cm(left)
    tf.margin_right = Cm(right)
    tf.margin_top = Cm(top)
    tf.margin_bottom = Cm(bottom)


# ---------------------------------------------------------------------------
# 各功能区调整
# ---------------------------------------------------------------------------


def _beautify_header(shapes) -> None:
    """顶部标题 + 副标题。"""
    title = shapes[3]
    sub = shapes[4]

    _set_run_size(title, 28)
    _set_geometry(title, top_cm=0.45, height_cm=1.10)
    _set_margins(title, top=0.0, bottom=0.0)
    _vertical_center(title)

    _set_run_size(sub, 13)
    _set_geometry(sub, top_cm=1.55, height_cm=0.65)
    _set_margins(sub, top=0.0, bottom=0.0)
    _vertical_center(sub)


def _beautify_top_signal_cards(shapes) -> None:
    """顶部三个圆角矩形：市场/客户/供应 信号。

    父卡片 top=4.37 height=1.83，新文字布局：
      - 标题 18pt @ 4.55cm 高 0.70cm
      - 副标题 12pt @ 5.30cm 高 0.55cm
    """
    top_title_idx = [7, 10, 13]
    top_sub_idx = [8, 11, 14]

    # 父卡片 width=7.49cm，标题原文本框只有 4.06 cm，18pt 大字会换行；
    # 把标题文本框宽度对齐到副标题(6.25cm)，且 left 对齐 2.34cm。
    for idx in top_title_idx:
        sh = shapes[idx]
        _set_run_size(sh, 18)
        _set_geometry(sh, top_cm=4.55, height_cm=0.70, width_cm=6.25)
        _set_margins(sh, top=0.0, bottom=0.0)
        _vertical_center(sh)

    for idx in top_sub_idx:
        sh = shapes[idx]
        _set_run_size(sh, 12)
        _set_geometry(sh, top_cm=5.30, height_cm=0.65)
        _set_margins(sh, top=0.0, bottom=0.0)
        _vertical_center(sh)


def _beautify_pivot_labels(shapes) -> None:
    """业务输入 / 一期功能(1-7) 两个枢纽标签。"""
    biz_input = shapes[20]
    phase_one = shapes[23]

    _set_run_size(biz_input, 12)
    _set_geometry(biz_input, top_cm=6.40, height_cm=0.50)
    _set_margins(biz_input, top=0.0, bottom=0.0)
    _vertical_center(biz_input)

    _set_run_size(phase_one, 13)
    _set_geometry(phase_one, top_cm=7.51, height_cm=0.50)
    _set_margins(phase_one, top=0.0, bottom=0.0)
    _vertical_center(phase_one)


def _beautify_function_cards(shapes) -> None:
    """7 个一期功能卡（含 1-7：情报/画像/供需/预测/定价/智能体/产销）。

    父卡片 top=8.18 height=3.51。子结构：
      - 顶部白底色块 top=8.18 height=0.86 (蓝色填充)
      - 顶部白字标题文本框：放大到 13pt，居中
      - 中间三栏特性：放大到 11pt
      - 底部小圆角 + 小标签：放大到 11pt
    """
    title_idx = [26, 33, 40, 47, 54, 61, 68]
    middle_idx = [27, 34, 41, 48, 55, 62, 69]
    badge_text_idx = [29, 36, 43, 50, 57, 64, 71]
    badge_box_idx = [28, 35, 42, 49, 56, 63, 70]

    for idx in title_idx:
        sh = shapes[idx]
        _set_run_size(sh, 13)
        _set_geometry(sh, top_cm=8.30, height_cm=0.62)
        _set_margins(sh, top=0.0, bottom=0.0)
        _vertical_center(sh)

    for idx in middle_idx:
        sh = shapes[idx]
        _set_run_size(sh, 11)
        # 三行内容均匀分布；line_spacing 控制紧凑度
        _set_geometry(sh, top_cm=9.25, height_cm=1.55)
        _set_margins(sh, top=0.0, bottom=0.0)
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                para.line_spacing = 1.2
        _vertical_center(sh)

    for idx in badge_box_idx:
        sh = shapes[idx]
        _set_geometry(sh, top_cm=10.92, height_cm=0.70)

    for idx in badge_text_idx:
        sh = shapes[idx]
        _set_run_size(sh, 11)
        _set_geometry(sh, top_cm=11.00, height_cm=0.55)
        _set_margins(sh, top=0.0, bottom=0.0)
        _vertical_center(sh)


def _beautify_review_section(shapes) -> None:
    """复盘校验维度 + D1~D5 卡。"""
    title_idx = 77
    sh = shapes[title_idx]
    _set_run_size(sh, 13)
    _set_geometry(sh, top_cm=11.85, height_cm=0.55)
    _set_margins(sh, top=0.0, bottom=0.0)
    _vertical_center(sh)

    # 父卡片 top=12.40 height=1.27
    # D 圆形 top=12.65 height=0.56
    groups = [
        (80, 81, 82, 83),   # D1 圆形 / D1 文字 / 标题 / 副标
        (85, 86, 87, 88),   # D2
        (90, 91, 92, 93),   # D3
        (95, 96, 97, 98),   # D4
        (100, 101, 102, 103),  # D5
    ]

    for circle_idx, dno_idx, title_idx_, sub_idx in groups:
        circle = shapes[circle_idx]
        _set_geometry(circle, top_cm=12.70, height_cm=0.62)

        dno = shapes[dno_idx]
        _set_run_size(dno, 12)
        _set_geometry(dno, top_cm=12.78, height_cm=0.45)
        _set_margins(dno, top=0.0, bottom=0.0)
        _vertical_center(dno)

        title = shapes[title_idx_]
        _set_run_size(title, 12.5)
        _set_geometry(title, top_cm=12.55, height_cm=0.50)
        _set_margins(title, top=0.0, bottom=0.0)
        _vertical_center(title)

        sub = shapes[sub_idx]
        _set_run_size(sub, 10.5)
        _set_geometry(sub, top_cm=13.10, height_cm=0.45)
        _set_margins(sub, top=0.0, bottom=0.0)
        _vertical_center(sub)


def _beautify_summary_card(shapes) -> None:
    """8 销售复盘 大卡片 + 校验/归因/优化 三个子项。

    父卡 top=13.97 height=3.20。顶部黄底 top=13.97 height=0.97。
    """
    big_title = shapes[106]   # "8 销售复盘"
    big_sub = shapes[107]      # 副标题
    # 原宽度 3.25 cm 装不下 18pt 5 字，扩到 5.0 cm；副标题相应右移、收窄。
    _set_run_size(big_title, 18)
    _set_geometry(big_title, top_cm=14.04, height_cm=0.75, width_cm=5.0)
    _set_margins(big_title, top=0.0, bottom=0.0)
    _vertical_center(big_title)

    _set_run_size(big_sub, 13)
    _set_geometry(big_sub, top_cm=14.13, height_cm=0.65, left_cm=7.40, width_cm=19.40)
    _set_margins(big_sub, top=0.0, bottom=0.0)
    _vertical_center(big_sub)

    # 三个子卡：校验/归因/优化
    # (key_box_idx, key_text_idx, sub_text_idx)
    sub_groups = [
        (109, 110, 111),
        (113, 114, 115),
        (117, 118, 119),
    ]
    for key_box_idx, key_text_idx, sub_text_idx in sub_groups:
        key_box = shapes[key_box_idx]
        _set_geometry(key_box, top_cm=15.55, height_cm=0.70)

        key_text = shapes[key_text_idx]
        _set_run_size(key_text, 13)
        _set_geometry(key_text, top_cm=15.63, height_cm=0.55)
        _set_margins(key_text, top=0.0, bottom=0.0)
        _vertical_center(key_text)

        sub = shapes[sub_text_idx]
        _set_run_size(sub, 12)
        _set_geometry(sub, top_cm=15.62, height_cm=0.55)
        _set_margins(sub, top=0.0, bottom=0.0)
        _vertical_center(sub)


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------


def beautify(src_path: Path, dst_path: Path) -> None:
    if not src_path.is_file():
        raise FileNotFoundError(f"找不到源文件: {src_path}")

    prs = Presentation(str(src_path))
    if not prs.slides:
        raise ValueError("PPT 中没有任何幻灯片")

    slide = prs.slides[0]
    shapes = list(slide.shapes)
    expected_count = 120
    if len(shapes) != expected_count:
        logger.warning("形状数量与预期不一致：当前 %d，预期 %d；脚本是按原版 shape index 写死的，"
                       "如改动过结构请先重新核对索引。", len(shapes), expected_count)

    _beautify_header(shapes)
    _beautify_top_signal_cards(shapes)
    _beautify_pivot_labels(shapes)
    _beautify_function_cards(shapes)
    _beautify_review_section(shapes)
    _beautify_summary_card(shapes)

    dst_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst_path))
    logger.info("已保存美化版到: %s", dst_path)


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--src", required=True, help="源 PPT 路径")
    parser.add_argument("--dst", required=True, help="输出 PPT 路径")
    return parser.parse_args()


def main() -> None:
    logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(message)s")
    args = _parse_args()
    beautify(Path(args.src), Path(args.dst))


if __name__ == "__main__":
    main()
