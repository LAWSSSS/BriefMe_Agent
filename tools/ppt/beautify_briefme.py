"""为已存在的 BriefMe 汇报 PPT 做"非破坏性美化"。

策略（关键）：
  - 不删除/不移动 任何已有形状（图片 / 文本框 / 占位符）
  - 只在底层（z-order 最低）追加装饰：色带、卡片、几何点缀
  - 在顶部 0~0.40 与底部 7.20~7.50 的"边缘空白带"上加品牌条 / 页脚
  - 给每页标题左侧追加一个红色色块，作为强调条

输出与输入保持等价，但视觉统一度大幅提升。
"""
from __future__ import annotations

import argparse
import logging
from pathlib import Path
from typing import Iterable, Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# 视觉规范（与原 PPT 主色 #16448C 对齐）
# ---------------------------------------------------------------------------
BRAND_DEEP = RGBColor(0x16, 0x44, 0x8C)  # 主深蓝（与原 PPT 完全一致）
BRAND_BLUE = RGBColor(0x2A, 0x6F, 0xC9)  # 中蓝
ACCENT_RED = RGBColor(0xC6, 0x00, 0x19)  # 装饰红
ACCENT_GOLD = RGBColor(0xE8, 0xA8, 0x2C)  # 强调金
GRAY_DARK = RGBColor(0x2D, 0x3A, 0x4A)
GRAY_MID = RGBColor(0x6B, 0x77, 0x83)
GRAY_LIGHT = RGBColor(0xC9, 0xD3, 0xDC)
BG_LIGHT = RGBColor(0xF3, 0xF7, 0xFB)
BG_INFO = RGBColor(0xE8, 0xF1, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_FAMILY = "微软雅黑"

# 各页的章节简称（用于页脚中段）
PAGE_SECTIONS = {
    1: "封面",
    2: "一. 背景",
    3: "一. 背景 · 镔鑫",
    4: "一. 背景 · 盛隆",
    5: "一. 背景 · 永锋",
    6: "二. 实现过程",
    7: "二. 实现过程",
    8: "二. 实现过程",
    9: "三. 应用效果 · 镔鑫",
    10: "三. 应用效果 · 盛隆",
    11: "三. 应用效果 · 永锋",
    12: "四. 未来预期",
}


# ---------------------------------------------------------------------------
# OOXML 工具：z-order 调整
# ---------------------------------------------------------------------------
def _send_to_back(shape) -> None:
    """把 shape 移到 spTree 中第一个内容形状之前（最底层）。

    spTree 子元素结构：nvGrpSpPr, grpSpPr, [sp/pic/grpSp/cxnSp/graphicFrame ...]。
    我们把目标 shape 元素移动到第一个"内容形状"之前。
    """
    el = shape._element
    parent = el.getparent()
    parent.remove(el)
    content_tags = {"sp", "pic", "grpSp", "cxnSp", "graphicFrame"}
    first_content = None
    for child in parent:
        tag = etree.QName(child).localname
        if tag in content_tags:
            first_content = child
            break
    if first_content is not None:
        first_content.addprevious(el)
    else:
        parent.append(el)


# ---------------------------------------------------------------------------
# 形状构造
# ---------------------------------------------------------------------------
def _add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: Optional[RGBColor] = None,
    line: Optional[RGBColor] = None,
    line_width_pt: float = 0.5,
    rounded: bool = False,
    no_fill: bool = False,
):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shp_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if no_fill:
        shp.fill.background()
    elif fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width_pt)
    shp.shadow.inherit = False
    return shp


def _add_oval(
    slide, left: float, top: float, size: float, *, fill: RGBColor
):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_textbox(slide, left: float, top: float, width: float, height: float):
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )


def _set_text(
    paragraph,
    text: str,
    *,
    size: int = 11,
    bold: bool = False,
    color: RGBColor = GRAY_DARK,
    align=None,
    font_family: str = FONT_FAMILY,
) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_family
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align


# ---------------------------------------------------------------------------
# 通用：顶部品牌带 + 底部页脚
# ---------------------------------------------------------------------------
def _add_top_band(slide) -> None:
    """页面顶部 0~0.22 in 的深蓝细品牌带（轻量化，避免与原红色标题块冲突）。"""
    # 主带：0 ~ 0.20
    _add_rect(slide, 0, 0, 13.333, 0.20, fill=BRAND_DEEP, rounded=False)
    # 金色超细强调线：0.20 ~ 0.215
    _add_rect(slide, 0, 0.20, 13.333, 0.015, fill=ACCENT_GOLD, rounded=False)

    # 左：BriefMe 文字
    box = _add_textbox(slide, 0.4, 0.0, 6.0, 0.20)
    box.text_frame.margin_top = box.text_frame.margin_bottom = 0
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_text(p, "BriefMe", size=10, bold=True, color=WHITE)
    _set_text(
        p, "  ｜  工业视觉项目智能晨报 Agent",
        size=8, color=GRAY_LIGHT,
    )

    # 右：项目方
    box = _add_textbox(slide, 7.5, 0.0, 5.7, 0.20)
    box.text_frame.margin_top = box.text_frame.margin_bottom = 0
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _set_text(
        p, "中冶赛迪（重庆）信息技术有限公司",
        size=8, color=GRAY_LIGHT,
    )


def _add_footer(slide, page_no: int, total: int, section: str) -> None:
    """页面底部 7.18~7.45 in 的页脚：左品牌、中章节、右页码。"""
    # 顶部分隔细线
    _add_rect(slide, 0.4, 7.20, 12.53, 0.012, fill=GRAY_LIGHT, rounded=False)

    # 左：品牌
    box = _add_textbox(slide, 0.4, 7.24, 4.5, 0.22)
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_text(p, "BriefMe", size=9, bold=True, color=BRAND_DEEP)
    _set_text(p, "  ·  中冶赛迪", size=9, color=GRAY_MID)

    # 中：章节
    box = _add_textbox(slide, 4.5, 7.24, 4.33, 0.22)
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_text(p, section, size=9, color=GRAY_MID)

    # 右：页码（高亮当前页）
    box = _add_textbox(slide, 8.83, 7.24, 4.10, 0.22)
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _set_text(p, f"{page_no:02d}", size=11, bold=True, color=BRAND_DEEP)
    _set_text(p, f"  /  {total:02d}", size=9, color=GRAY_MID)


# ---------------------------------------------------------------------------
# 通用：标题左侧色块装饰
# ---------------------------------------------------------------------------
def _add_title_accent(slide) -> None:
    """在原标题（红底白字）下方加一条金色细线 + 末端金点，呼应顶部带。

    原标题位置（占位符 1）：L=0.97 T=0.52 W=7.76 H=0.44 → 底部 0.96
    （原标题已是红底白字，左侧不再加红色色块以避免色块叠加）。
    """
    # 标题底部金色细横线
    _add_rect(slide, 0.97, 0.99, 6.5, 0.018, fill=ACCENT_GOLD, rounded=False)
    # 末端深蓝小点（与顶部品牌色呼应）
    _add_oval(slide, 7.45, 0.97, 0.07, fill=BRAND_DEEP)


# ---------------------------------------------------------------------------
# 第 1 页（封面）：底层渐变 + 副标题 + 几何装饰
# ---------------------------------------------------------------------------
def _decorate_cover(slide) -> None:
    """封面页特殊处理：
    1. 底层加大色块装饰（左上深蓝、右下浅蓝）
    2. 加副标题"让算法工程师从重复劳动中解放出来"
    3. 加 4 个关键词标签
    """
    # 1) 底层装饰（这些会被 _send_to_back 推到最底层）
    bg_shapes = []
    # 左上深蓝大色块（让 logo 更醒目）
    bg_shapes.append(
        _add_rect(slide, -0.5, -0.5, 5.5, 2.5, fill=BRAND_DEEP, rounded=False)
    )
    # 右上柔和装饰圆（位于标题区域之外，远离右下 CISDigital 标识）
    bg_shapes.append(
        _add_oval(slide, 11.0, -1.5, 4.0, fill=BG_INFO)
    )
    # 左下小色块装饰（对称右上圆形）
    bg_shapes.append(
        _add_rect(slide, -0.3, 6.05, 1.2, 0.10, fill=ACCENT_GOLD, rounded=False)
    )
    # 中部柔和水平横线（标题下方分隔线）
    bg_shapes.append(
        _add_rect(slide, 0.5, 3.95, 12.33, 0.012, fill=GRAY_LIGHT, rounded=False)
    )
    # 主标题上方金红组合短线（增加品牌识别）
    bg_shapes.append(
        _add_rect(slide, 6.20, 1.95, 0.35, 0.06, fill=ACCENT_GOLD, rounded=False)
    )
    bg_shapes.append(
        _add_rect(slide, 6.65, 1.95, 0.50, 0.06, fill=ACCENT_RED, rounded=False)
    )

    # 把上述底层装饰送到最底层
    for shp in bg_shapes:
        _send_to_back(shp)

    # 2) 副标题（在主标题下方）—— 这个不需要送底层，正常显示
    box = _add_textbox(slide, 0.0, 3.50, 13.12, 0.55)
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_text(
        p, "让算法工程师从重复劳动中解放出来",
        size=18, bold=False, color=GRAY_DARK,
    )

    # 3) 关键词标签（4 个胶囊）
    keywords = [
        ("多系统汇聚", BRAND_BLUE),
        ("智能统计", ACCENT_GOLD),
        ("异常归因", ACCENT_RED),
        ("对话式查询", BRAND_DEEP),
    ]
    pill_w = 1.55
    pill_h = 0.36
    gap = 0.20
    total_w = pill_w * len(keywords) + gap * (len(keywords) - 1)
    start_x = (13.333 - total_w) / 2
    y = 6.10
    for label, color in keywords:
        # 胶囊背景：白底 + 彩色描边
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(start_x), Inches(y), Inches(pill_w), Inches(pill_h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = color
        shp.line.width = Pt(1.2)
        shp.shadow.inherit = False
        # 文字
        tb = _add_textbox(slide, start_x, y + 0.04, pill_w, pill_h - 0.08)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_text(p, label, size=12, bold=True, color=color)
        start_x += pill_w + gap


# ---------------------------------------------------------------------------
# 第 2 页：底层卡片 + 章节标识
# ---------------------------------------------------------------------------
def _decorate_page2_background(slide) -> None:
    """在第 2 页"背景"页加底层卡片，让散落的文字看起来分组明确。

    原页面布局解读：
      - 文本框 6（"随着...大量检测数据"）在 L=1.16 T=1.65 → 是引言段
      - 文本框 5（"模型检测结果历史记录..."）在 L=8.66 T=1.32 → 是数据类型 list
      - 文本框 8（"当前存在的问题"）+ 7（4 条问题） → 中部
      - 文本框 3（"多系统数据自动汇聚能力..."）在 L=4.08 T=4.54 → 解决方案
      - 文本框 9（产品名）+ 10/11（红色标语）→ 右下结论

    我们加 3 张底层卡片：
      上: 引言 + 数据类型清单 (T=1.10 ~ 2.55)
      中: 当前问题 (T=2.65 ~ 4.30)
      下: 解决方案 + 产品标语 (T=4.40 ~ 7.10)
    """
    bg_shapes = []
    # 上卡片（淡蓝）
    bg_shapes.append(
        _add_rect(
            slide, 0.7, 1.10, 11.93, 1.45,
            fill=BG_LIGHT, rounded=True,
        )
    )
    # 上卡片左侧色条
    bg_shapes.append(
        _add_rect(slide, 0.7, 1.10, 0.10, 1.45, fill=BRAND_DEEP, rounded=False)
    )

    # 中卡片（更暗一点）
    bg_shapes.append(
        _add_rect(
            slide, 0.7, 2.65, 11.93, 1.65,
            fill=RGBColor(0xFD, 0xF1, 0xEC), rounded=True,
        )
    )
    bg_shapes.append(
        _add_rect(slide, 0.7, 2.65, 0.10, 1.65, fill=ACCENT_RED, rounded=False)
    )

    # 下卡片（绿色调，代表方案）
    bg_shapes.append(
        _add_rect(
            slide, 0.7, 4.40, 11.93, 2.55,
            fill=BG_INFO, rounded=True,
        )
    )
    bg_shapes.append(
        _add_rect(slide, 0.7, 4.40, 0.10, 2.55, fill=BRAND_BLUE, rounded=False)
    )

    # 全部送至底层
    for shp in bg_shapes:
        _send_to_back(shp)

    # 在卡片左上角加章节小标签
    section_labels = [
        (1.12, 1.18, "现状", BRAND_DEEP),
        (1.12, 2.73, "痛点", ACCENT_RED),
        (1.12, 4.48, "方案", BRAND_BLUE),
    ]
    for x, y, label, color in section_labels:
        # 标签背景
        _add_rect(slide, x, y, 0.55, 0.26, fill=color, rounded=True)
        # 标签文字
        tb = _add_textbox(slide, x, y + 0.02, 0.55, 0.22)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_text(p, label, size=10, bold=True, color=WHITE)


# ---------------------------------------------------------------------------
# 第 12 页：右侧空白处加未来规划要点卡片
# ---------------------------------------------------------------------------
def _decorate_page12(slide) -> None:
    """第 12 页只有左侧一张图，右侧 8.3~13.0 大片空白，加 3 张未来规划卡片。"""
    items = [
        (
            "短期",
            "更多项目接入",
            "持续接入更多视觉项目（钢铁、汽车、新能源）的检测数据，覆盖更广业务场景。",
            BRAND_DEEP,
        ),
        (
            "中期",
            "多模态对话",
            "在 BriefMe 中集成图像/视频回看与现场采图溯源，让结果可追到原始证据。",
            BRAND_BLUE,
        ),
        (
            "远期",
            "决策辅助闭环",
            "结合大模型的分析，自动生成模型迭代建议、样本采集策略，形成『诊断-改进』闭环。",
            ACCENT_RED,
        ),
    ]
    left = 8.40
    width = 4.80
    top = 1.80
    height = 1.56  # 三张卡片 + 间距 = (3 * 1.56) + (2 * 0.20) = 5.08，5.08 < 5.30 OK
    gap = 0.20

    # 区块标题
    box = _add_textbox(slide, left, 1.30, width, 0.45)
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_text(p, "未来规划路线", size=18, bold=True, color=BRAND_DEEP)

    for stage, title, desc, color in items:
        # 卡片
        _add_rect(
            slide, left, top, width, height,
            fill=BG_LIGHT, line=GRAY_LIGHT, line_width_pt=0.5, rounded=True,
        )
        # 左侧色条
        _add_rect(slide, left, top, 0.10, height, fill=color, rounded=False)

        # 阶段徽章
        _add_rect(slide, left + 0.25, top + 0.12, 0.50, 0.28, fill=color, rounded=True)
        tb = _add_textbox(slide, left + 0.25, top + 0.14, 0.50, 0.25)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_text(p, stage, size=10, bold=True, color=WHITE)

        # 主标题
        tb = _add_textbox(slide, left + 0.85, top + 0.10, width - 1.0, 0.36)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _set_text(p, title, size=14, bold=True, color=BRAND_DEEP)

        # 描述
        tb = _add_textbox(slide, left + 0.25, top + 0.50, width - 0.45, height - 0.55)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _set_text(p, desc, size=10, color=GRAY_DARK)

        top += height + gap


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------
def beautify(input_path: Path, output_path: Path) -> Path:
    input_path = Path(input_path).resolve()
    output_path = Path(output_path).resolve()

    if not input_path.exists():
        raise FileNotFoundError(f"输入 PPT 不存在：{input_path}")

    prs = Presentation(str(input_path))
    total_pages = len(prs.slides)
    logger.info("开始美化 %s（共 %d 页）", input_path, total_pages)

    for idx, slide in enumerate(prs.slides):
        page_no = idx + 1
        section = PAGE_SECTIONS.get(page_no, f"第 {page_no} 页")

        # 第 1 页（封面）：仅做底层装饰 + 关键词，不加顶部带 / 页脚 / 标题色块
        if page_no == 1:
            _decorate_cover(slide)
            continue

        # 其他页：通用顶部带 + 页脚 + 标题色块
        _add_top_band(slide)
        _add_footer(slide, page_no, total_pages, section)
        _add_title_accent(slide)

        # 各页特殊处理
        if page_no == 2:
            _decorate_page2_background(slide)
        elif page_no == 12:
            _decorate_page12(slide)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info("美化完成 → %s", output_path)
    return output_path


def _build_argparser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="对 BriefMe 汇报 PPT 做非破坏性美化"
    )
    parser.add_argument(
        "--input",
        default="演示文件/briefme_agent_ppt.original.pptx",
        help="输入 PPT 路径（默认：演示文件/briefme_agent_ppt.original.pptx）",
    )
    parser.add_argument(
        "--output",
        default="演示文件/briefme_agent_ppt.pptx",
        help="输出 PPT 路径（默认：演示文件/briefme_agent_ppt.pptx，会覆盖）",
    )
    return parser


def main() -> None:
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s | %(levelname)s | %(message)s",
    )
    args = _build_argparser().parse_args()
    beautify(Path(args.input), Path(args.output))


if __name__ == "__main__":
    main()
