"""镔鑫 PPT 生成 smoke 测试。

走 ``agent.scrap.ppt_writer.write_stats_pptx``，主路径是自研 builder
(``agent.scrap.ppt_builder.build_binxin_ppt``)，不依赖外部 API。

运行：
  /opt/anaconda3/bin/python -m pytest tests/test_scrap_ppt_unit.py -q
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from agent.scrap.models import DailyScrapStats, MaterialEntry, TruckStat
from agent.scrap.ppt_writer import (
    PPTGenerationError,
    write_stats_pptx,
)


def _make_stats(
    date_str: str,
    acc: float,
    err: float,
    wd: float,
    wr: float,
    *,
    n_errors: int = 0,
) -> DailyScrapStats:
    s = DailyScrapStats(date=date_str)
    s.total_trucks = 10
    s.eligible_trucks = 10
    s.main_same_count = int(round(acc / 10))
    s.avg_error_rate = err
    s.avg_weight_diff = wd
    s.avg_weight_ratio = wr
    # 模拟错判车：按需要插入 n_errors 辆
    for i in range(n_errors):
        t = TruckStat(
            date=date_str,
            car_number=f"桂P{50000 + i}",
            station_number=53,
            flow_code=f"fc-{date_str}-{i}",
        )
        t.manual_main = MaterialEntry(steel_type=4, steel_level=0, rate=70.0)
        t.ai_main = MaterialEntry(
            steel_type=3, steel_level=0, rate=70.0 - (i + 1) * 5.0
        )
        t.is_eligible = True
        t.main_same = False
        t.diff_rate = (i + 1) * 5.0
        s.trucks.append(t)
    return s


def test_ppt_happy_path(tmp_path):
    """6 天有效数据 → 生成单页 PPT，含 chart + 元信息卡 + KPI 卡 + 改进建议"""
    stats_list = [
        _make_stats("2026-04-14", 60.0, 12.5, 80.0, 0.95, n_errors=4),
        _make_stats("2026-04-15", 67.0, 11.2, 75.0, 1.05, n_errors=3),
        _make_stats("2026-04-16", 73.0, 10.0, 70.0, 1.10, n_errors=3),
        _make_stats("2026-04-17", 80.0, 8.5, 65.0, 1.08, n_errors=2),
        _make_stats("2026-04-18", 86.0, 7.0, 58.0, 1.02, n_errors=1),
        _make_stats("2026-04-19", 90.0, 5.5, 50.0, 0.98, n_errors=1),
    ]
    out = tmp_path / "test.pptx"
    result = write_stats_pptx(
        stats_list, out,
        start_date="2026-04-14",
        end_date="2026-04-19",
        target_pct=95.0,
    )
    assert result.exists()
    assert result.stat().st_size > 10000

    from pptx import Presentation

    prs = Presentation(str(result))
    assert len(prs.slides) == 1
    slide = prs.slides[0]

    # 1. 1 张可编辑图表（折线 + 目标线）
    chart_shapes = [shape for shape in slide.shapes if shape.has_chart]
    assert len(chart_shapes) == 1, "应该有 1 个折线图（识别率趋势 + 目标线）"
    chart = chart_shapes[0].chart
    assert len(list(chart.series)) == 2, "图表应该有 2 个系列：识别率 + 目标线"

    # 2. 1 张表格（错判 Top5）
    table_shapes = [shape for shape in slide.shapes if shape.has_table]
    assert len(table_shapes) == 1, "应该有 1 个错判 Top5 表格"

    # 3. 关键文本（标题/元信息卡/KPI/改进建议）
    text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
    text_blob = "\n".join(s.text_frame.text for s in text_shapes)
    for keyword in (
        "镔鑫废钢检判",
        "数据来源",
        "统计周期",
        "有效车次",
        "错判车次",
        "周期核心 KPI",
        "改进建议",
        "BriefMe",
        "中冶赛迪",
    ):
        assert keyword in text_blob, f"PPT 应该包含关键字「{keyword}」"


def test_ppt_no_errors_shows_empty_table_state(tmp_path):
    """全部判断正确（n_errors=0）→ 不出错判表，但有"无错判车次"提示"""
    stats_list = [
        _make_stats("2026-04-14", 100.0, 5.0, 50.0, 1.0, n_errors=0),
        _make_stats("2026-04-15", 100.0, 5.0, 50.0, 1.0, n_errors=0),
    ]
    out = tmp_path / "no_errors.pptx"
    write_stats_pptx(
        stats_list, out, start_date="2026-04-14", end_date="2026-04-15"
    )
    from pptx import Presentation

    prs = Presentation(str(out))
    text_blob = "\n".join(
        s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame
    )
    # 无错判时，应该显示"恭喜，全部主料判定一致"提示
    assert "无错判车次" in text_blob


def test_ppt_empty_stats_raises(tmp_path):
    """空列表 → 抛 PPTGenerationError"""
    with pytest.raises(PPTGenerationError):
        write_stats_pptx([], tmp_path / "x.pptx")


def test_ppt_insufficient_rows_raises(tmp_path):
    """只有 1 天有效数据 → 抛 PPTGenerationError（趋势图至少 2 行）"""
    stats_list = [_make_stats("2026-04-14", 70.0, 10.0, 80.0, 1.0)]
    with pytest.raises(PPTGenerationError, match="不足"):
        write_stats_pptx(stats_list, tmp_path / "x.pptx")


def test_ppt_skips_zero_eligible(tmp_path):
    """eligible_trucks=0 的天会被跳过，剩 2 天数据仍能出图"""
    s_skip = DailyScrapStats(date="2026-04-13")
    s_skip.total_trucks = 0
    s_skip.eligible_trucks = 0
    stats_list = [
        s_skip,
        _make_stats("2026-04-14", 70.0, 10.0, 80.0, 1.0),
        _make_stats("2026-04-15", 75.0, 9.0, 75.0, 1.05),
    ]
    out = tmp_path / "skip.pptx"
    write_stats_pptx(stats_list, out)
    assert out.exists() and out.stat().st_size > 10000


def test_ppt_legacy_fallback_still_works(tmp_path):
    """显式 use_legacy=True 仍走同事 skill（保留兼容）"""
    stats_list = [
        _make_stats("2026-04-14", 70.0, 10.0, 80.0, 1.0),
        _make_stats("2026-04-15", 75.0, 9.0, 75.0, 1.05),
        _make_stats("2026-04-16", 80.0, 8.0, 70.0, 1.08),
    ]
    out = tmp_path / "legacy.pptx"
    result = write_stats_pptx(stats_list, out, use_legacy=True)
    assert result.exists() and result.stat().st_size > 10000
